import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import openpyxl
from typing import List


def parse_document(file_path: str, file_type: str) -> str:
    """Extract raw text from PDF, PPTX, DOCX, or XLSX."""
    file_type = file_type.lower()
    if file_type == "pdf":
        return _parse_pdf(file_path)
    elif file_type == "pptx":
        return _parse_pptx(file_path)
    elif file_type == "docx":
        return _parse_docx(file_path)
    elif file_type == "xlsx":
        return _parse_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def _parse_pdf(path: str) -> str:
    doc = fitz.open(path)
    return "\n\n".join(page.get_text() for page in doc)


def _parse_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        slide_text = " ".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
        )
        if slide_text:
            texts.append(slide_text)
    return "\n\n".join(texts)


def _parse_docx(path: str) -> str:
    doc = Document(path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    rows = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                rows.append(row_text)
    return "\n".join(rows)


def chunk_text(text: str, chunk_size: int = 2000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks, preferring paragraph boundaries."""
    paragraphs = text.split("\n\n")
    chunks = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) <= chunk_size:
            current += para + "\n\n"
        else:
            if current.strip():
                chunks.append(current.strip())
            current = current[-overlap:] + para + "\n\n"
    if current.strip():
        chunks.append(current.strip())
    return chunks
